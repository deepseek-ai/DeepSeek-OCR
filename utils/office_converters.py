"""
Microsoft Office Format Converters
Convert DOCX, PPTX, XLSX to images for OCR processing
"""

import io
import tempfile
from pathlib import Path
from typing import List, Optional
from PIL import Image


class OfficeConverter:
    """Base class for office document converters"""

    @staticmethod
    def convert_to_images(file_bytes: bytes, file_type: str, dpi: int = 150) -> List[Image.Image]:
        """Convert office document to images"""
        if file_type == 'docx':
            return DOCXConverter.convert(file_bytes, dpi)
        elif file_type == 'pptx':
            return PPTXConverter.convert(file_bytes, dpi)
        elif file_type == 'xlsx':
            return XLSXConverter.convert(file_bytes, dpi)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")


class DOCXConverter:
    """Convert DOCX documents to images"""

    @staticmethod
    def convert(file_bytes: bytes, dpi: int = 150) -> List[Image.Image]:
        """Convert DOCX to images using docx2pdf or alternative methods"""
        try:
            # Method 1: Try using pypandoc with wkhtmltopdf
            return DOCXConverter._convert_via_pandoc(file_bytes, dpi)
        except:
            try:
                # Method 2: Try using python-docx to extract content and render
                return DOCXConverter._convert_via_docx(file_bytes, dpi)
            except Exception as e:
                raise Exception(f"Failed to convert DOCX: {str(e)}")

    @staticmethod
    def _convert_via_docx(file_bytes: bytes, dpi: int = 150) -> List[Image.Image]:
        """Convert DOCX by extracting content and rendering to image"""
        from docx import Document
        from PIL import Image, ImageDraw, ImageFont
        import textwrap

        doc = Document(io.BytesIO(file_bytes))

        images = []
        page_width = int(8.5 * dpi)  # Letter size
        page_height = int(11 * dpi)

        current_y = 50
        img = Image.new('RGB', (page_width, page_height), 'white')
        draw = ImageDraw.Draw(img)

        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
            font_bold = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
        except:
            font = ImageFont.load_default()
            font_bold = font

        margin = 100

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                current_y += 20
                continue

            # Check if heading
            is_heading = para.style.name.startswith('Heading')
            current_font = font_bold if is_heading else font

            # Wrap text
            wrapped_lines = textwrap.wrap(text, width=80)

            for line in wrapped_lines:
                if current_y > page_height - 100:
                    # New page
                    images.append(img)
                    img = Image.new('RGB', (page_width, page_height), 'white')
                    draw = ImageDraw.Draw(img)
                    current_y = 50

                draw.text((margin, current_y), line, fill='black', font=current_font)
                current_y += 35 if is_heading else 30

        if current_y > 50:  # Add last page if it has content
            images.append(img)

        return images if images else [Image.new('RGB', (page_width, page_height), 'white')]

    @staticmethod
    def _convert_via_pandoc(file_bytes: bytes, dpi: int = 150) -> List[Image.Image]:
        """Convert DOCX via pandoc to HTML then render"""
        import pypandoc
        import tempfile
        import os

        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_docx:
            tmp_docx.write(file_bytes)
            docx_path = tmp_docx.name

        try:
            # Convert to HTML
            html_output = pypandoc.convert_file(docx_path, 'html')

            # Render HTML to image (simplified)
            # This would require additional libraries like selenium or weasyprint
            # For now, fall back to docx method
            raise NotImplementedError("Pandoc conversion needs additional setup")

        finally:
            os.unlink(docx_path)


class PPTXConverter:
    """Convert PPTX presentations to images"""

    @staticmethod
    def convert(file_bytes: bytes, dpi: int = 150) -> List[Image.Image]:
        """Convert PPTX slides to images"""
        try:
            from pptx import Presentation
            from PIL import Image, ImageDraw, ImageFont
            import io

            prs = Presentation(io.BytesIO(file_bytes))
            images = []

            # Slide dimensions (16:9 aspect ratio)
            slide_width = int(10 * dpi)  # 10 inches
            slide_height = int(10 * 9 / 16 * dpi)

            for slide_num, slide in enumerate(prs.slides):
                img = Image.new('RGB', (slide_width, slide_height), 'white')
                draw = ImageDraw.Draw(img)

                try:
                    font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
                    font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
                except:
                    font = ImageFont.load_default()
                    font_title = font

                current_y = 100

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()

                        # Determine if title (usually first text box)
                        is_title = current_y < 200
                        current_font = font_title if is_title else font

                        # Draw text
                        lines = text.split('\n')
                        for line in lines:
                            if current_y < slide_height - 100:
                                draw.text((150, current_y), line, fill='black', font=current_font)
                                current_y += 60 if is_title else 40

                        current_y += 40  # Space between text boxes

                    # Try to extract images from slide
                    if hasattr(shape, "image"):
                        try:
                            image_stream = io.BytesIO(shape.image.blob)
                            embedded_img = Image.open(image_stream)

                            # Resize to fit
                            max_width = slide_width - 300
                            max_height = slide_height - current_y - 100

                            embedded_img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)

                            # Paste onto slide
                            x_pos = (slide_width - embedded_img.width) // 2
                            img.paste(embedded_img, (x_pos, current_y))
                            current_y += embedded_img.height + 40
                        except:
                            pass

                images.append(img)

            return images if images else [Image.new('RGB', (slide_width, slide_height), 'white')]

        except Exception as e:
            raise Exception(f"Failed to convert PPTX: {str(e)}")


class XLSXConverter:
    """Convert XLSX spreadsheets to images"""

    @staticmethod
    def convert(file_bytes: bytes, dpi: int = 150) -> List[Image.Image]:
        """Convert XLSX sheets to images"""
        try:
            import openpyxl
            from PIL import Image, ImageDraw, ImageFont
            import io

            wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
            images = []

            # Page dimensions
            page_width = int(11 * dpi)  # Landscape letter
            page_height = int(8.5 * dpi)

            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf", 20)
            except:
                font = ImageFont.load_default()

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]

                img = Image.new('RGB', (page_width, page_height), 'white')
                draw = ImageDraw.Draw(img)

                # Draw sheet name
                draw.text((50, 30), f"Sheet: {sheet_name}", fill='black', font=font)

                current_y = 100
                cell_height = 30
                cell_padding = 10

                # Get active area
                max_row = min(ws.max_row, 50)  # Limit to 50 rows
                max_col = min(ws.max_column, 15)  # Limit to 15 columns

                # Calculate column widths
                col_widths = []
                for col in range(1, max_col + 1):
                    max_width = 100
                    for row in range(1, max_row + 1):
                        cell = ws.cell(row=row, column=col)
                        if cell.value:
                            cell_text = str(cell.value)
                            bbox = draw.textbbox((0, 0), cell_text, font=font)
                            text_width = bbox[2] - bbox[0]
                            max_width = max(max_width, text_width + 20)
                    col_widths.append(min(max_width, 200))

                # Draw cells
                for row in range(1, max_row + 1):
                    if current_y > page_height - 100:
                        # New page
                        images.append(img)
                        img = Image.new('RGB', (page_width, page_height), 'white')
                        draw = ImageDraw.Draw(img)
                        current_y = 50

                    current_x = 50

                    for col in range(1, max_col + 1):
                        cell = ws.cell(row=row, column=col)
                        cell_width = col_widths[col - 1]

                        # Draw cell border
                        draw.rectangle(
                            [(current_x, current_y), (current_x + cell_width, current_y + cell_height)],
                            outline='black',
                            fill='lightgray' if row == 1 else 'white'
                        )

                        # Draw cell text
                        if cell.value:
                            cell_text = str(cell.value)[:30]  # Truncate long text
                            draw.text(
                                (current_x + cell_padding, current_y + cell_padding // 2),
                                cell_text,
                                fill='black',
                                font=font
                            )

                        current_x += cell_width

                    current_y += cell_height

                images.append(img)

            return images if images else [Image.new('RGB', (page_width, page_height), 'white')]

        except Exception as e:
            raise Exception(f"Failed to convert XLSX: {str(e)}")
